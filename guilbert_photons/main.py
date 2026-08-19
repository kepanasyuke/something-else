import asyncio
import io
import base64
import json
import logging
import os
from pathlib import Path
import numpy as np
import matplotlib
# Инициализируем безэкранный бэкенд для работы matplotlib в потоках FastAPI
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from contextlib import asynccontextmanager
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pptx import Presentation
from pptx.util import Inches
from pydantic import BaseModel, Field
from typing import List, Optional

# --- СХЕМЫ ДАННЫХ КОНТРАКТА ---
class DocumentSchema(BaseModel):
    id: int
    rubrics: List[str] = Field(default_factory=list)
    title: Optional[str] = None
    text: str
    created_date: Optional[str] = None
    journal: Optional[str] = None
    authors: List[str] = Field(default_factory=list)
    doi: Optional[str] = None
    source_url: Optional[str] = None
    is_open_access: bool = False

class SearchRequest(BaseModel):
    query: str = Field(min_length=1, max_length=200)
    limit: int = Field(default=30, ge=1, le=100)
    offset: int = Field(default=0, ge=0)

class SearchResponse(BaseModel):
    results: List[DocumentSchema]
    total: int

# Новая Pydantic-модель для реактивного управления графиками
class PlotRequest(BaseModel):
    frequency: float = Field(ge=0.1, le=10.0)

class PlotResponse(BaseModel):
    packet_plot: str  # Картинка 1 в формате Base64
    area_plot: str    # Картинка 2 в формате Base64
    scale_plot: str   # Картинка 3 в формате Base64

class PowerPointRequest(BaseModel):
    frequency: float
    intensity: str
    amplitudes: List[float] = Field(default_factory=list)

KNOWLEDGE_BASE_PATH = Path(__file__).with_name("knowledge_base.json")


def load_knowledge_base() -> List[DocumentSchema]:
    with KNOWLEDGE_BASE_PATH.open(encoding="utf-8") as source:
        return [DocumentSchema.model_validate(document) for document in json.load(source)]


knowledge_base = load_knowledge_base()

# Настройка Enterprise-логирования
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S"
)
logger = logging.getLogger("QuantumEngine")

@asynccontextmanager
async def lifespan(app: FastAPI):
    logger.info("--- [Инициализация] Математический движок Гильбертовых полей запущен ---")
    yield
    logger.info("--- [Остановка] Завершение работы вычислительных ядер ---")

app = FastAPI(title="Search & Quantum Field Service", version="3.0", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=os.getenv(
        "CORS_ORIGINS",
        "http://localhost:8000,http://127.0.0.1:8000",
    ).split(","),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- 🔬 МАТЕМАТИЧЕСКОЕ ЯДРО (ОБСЧЕТ И ОТРИСОВКА ЧЕРТЕЖЕЙ) ---
PLOT_BACKGROUND = "#08111a"


def encode_figure(figure: object, background: str) -> str:
    buffer = io.BytesIO()
    try:
        figure.savefig(buffer, format="png", dpi=120, bbox_inches="tight", facecolor=background)
        buffer.seek(0)
        encoded = base64.b64encode(buffer.read()).decode("utf-8")
    finally:
        plt.close(figure)
        buffer.close()
    return encoded


def render_packet_plot(omega: float) -> str:
    fig, ax = plt.subplots(figsize=(5, 3.5), facecolor=PLOT_BACKGROUND)
    ax.set_facecolor(PLOT_BACKGROUND)
    x = np.linspace(-5, 5, 400)
    sigma = 1.0
    k0 = omega * 2.0
    psi_real = np.exp(-x**2 / (4 * sigma**2)) * np.cos(k0 * x)
    probability_density = np.exp(-x**2 / (2 * sigma**2))
    
    ax.plot(x, psi_real, color="#62d9d0", linewidth=1.8, label="Real Ψ(x)")
    ax.plot(x, probability_density, color="#f4b860", linewidth=1.2, linestyle="--", label="|Ψ|² Density")
    ax.set_title(r"$\Psi(x) = e^{-\frac{x^2}{4\sigma^2}} \cdot e^{ik_0 x}$", fontsize=11, color="#dce9e8")
    ax.set_ylim(-1.2, 1.2)
    ax.axis('off')
    return encode_figure(fig, PLOT_BACKGROUND)


def render_area_plot(omega: float) -> str:
    fig, ax = plt.subplots(figsize=(5, 3.5), facecolor=PLOT_BACKGROUND)
    ax.set_facecolor(PLOT_BACKGROUND)
    random_generator = np.random.default_rng(int(omega * 10))
    num_nodes = int(omega * 4) + 4
    nodes_x = random_generator.uniform(-2, 2, num_nodes)
    nodes_y = random_generator.uniform(-2, 2, num_nodes)
    for i in range(num_nodes):
        for j in range(i+1, num_nodes):
            if np.hypot(nodes_x[i]-nodes_x[j], nodes_y[i]-nodes_y[j]) < 1.8:
                ax.plot([nodes_x[i], nodes_x[j]], [nodes_y[i], nodes_y[j]], color="#62d9d0", linewidth=0.5, alpha=0.3)
                
    ax.scatter(nodes_x, nodes_y, color="#f4b860", s=45, zorder=5, edgecolors="#ffffff", linewidths=0.5)
    ax.set_title(r"$A = 8\pi\gamma \ell_P^2 \sum \sqrt{j_i(j_i+1)}$", fontsize=11, color="#dce9e8")
    ax.axis('off')
    return encode_figure(fig, PLOT_BACKGROUND)


def render_scale_plot(omega: float) -> str:
    fig, ax = plt.subplots(figsize=(5, 3.5), facecolor=PLOT_BACKGROUND)
    ax.set_facecolor(PLOT_BACKGROUND)
    mu = np.linspace(0.1, 10, 200)
    alpha = 1.0 / (0.5 * omega * np.log(mu + 1) + 0.2)
    ax.plot(mu, alpha, color="#62d9d0", linewidth=2)
    ax.axvline(x=7.5, color="#f06f61", linestyle=":", linewidth=1.2, label="Planck Scale")
    ax.set_title(r"$\frac{d\alpha}{d\ln\mu} = \beta(\alpha)$ (Field Scale)", fontsize=11, color="#dce9e8")
    ax.axis('off')
    return encode_figure(fig, PLOT_BACKGROUND)


def compute_and_render_plots(omega: float) -> dict:
    with plt.style.context("dark_background"):
        plots = {
            "packet": render_packet_plot(omega),
            "area": render_area_plot(omega),
            "scale": render_scale_plot(omega),
        }
    return plots

# --- 🚀 REST ЭНДПОИНТЫ ---

@app.get("/", response_class=FileResponse)
async def root():
    return FileResponse(Path(__file__).with_name("guilbert.html"), media_type="text/html")

@app.get("/guilbert.css", response_class=FileResponse)
async def stylesheet():
    return FileResponse(Path(__file__).with_name("guilbert.css"), media_type="text/css")

def build_powerpoint(request: PowerPointRequest) -> io.BytesIO:
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Quantum Field Report"
    title_slide.placeholders[1].text = (
        f"Частота ω: {request.frequency:.1f}\n"
        f"Режим: {request.intensity}"
    )

    data_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    data_slide.shapes.title.text = "Амплитуды волнового поля"
    textbox = data_slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(4.5)
    )
    textbox.text_frame.text = "\n".join(
        f"{index + 1}. {amplitude:.4f}"
        for index, amplitude in enumerate(request.amplitudes)
    ) or "Данные отсутствуют"

    output = io.BytesIO()
    presentation.save(output)
    output.seek(0)
    return output


@app.post("/api/v1/generate-powerpoint")
async def generate_powerpoint(request: PowerPointRequest):
    output = build_powerpoint(request)
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=quantum_report.pptx"},
    )

@app.post("/generate-plots", response_model=PlotResponse)
async def generate_quantum_plots(request: PlotRequest):
    logger.info("Запрос на рендеринг квантового поля. Координата частоты ω = %s", request.frequency)
    try:
        # [Архитектурный паттерн]: выносим тяжелую графику Matplotlib в thread pool, 
        # защищая Event Loop сервера от зависания
        loop = asyncio.get_running_loop()
        encoded_plots = await loop.run_in_executor(None, compute_and_render_plots, request.frequency)
        
        return PlotResponse(
            packet_plot=encoded_plots['packet'],
            area_plot=encoded_plots['area'],
            scale_plot=encoded_plots['scale']
        )
    except Exception as e:
        logger.error("Критический сбой математического ядра: %s", e)
        raise HTTPException(status_code=500, detail="Ошибка вычисления топологии поля")

@app.post("/search", response_model=SearchResponse)
async def search(request: SearchRequest):
    logger.info("Локальный поиск мини-вики. Запрос: '%s'", request.query)
    query = request.query.casefold()
    matching_documents = [
        document for document in knowledge_base
        if not query or any(
            query in value.casefold()
            for value in (
                document.text,
                document.title or "",
                document.journal or "",
                *document.rubrics,
                *document.authors,
            )
        )
    ]
    page = matching_documents[request.offset:request.offset + request.limit]
    return SearchResponse(results=page, total=len(matching_documents))

@app.delete("/documents/{doc_id}")
async def delete_doc(doc_id: int):
    global knowledge_base
    knowledge_base = [document for document in knowledge_base if document.id != doc_id]
    return {"status": "deleted"}

@app.get("/health", include_in_schema=False)
async def health_check():
    return {"status": "ok"}




